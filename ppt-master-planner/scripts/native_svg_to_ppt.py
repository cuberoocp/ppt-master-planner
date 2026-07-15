#!/usr/bin/env python3
"""SVG to Native PPTX Converter v2.0

Converts hand-crafted SVG pages into native PowerPoint shapes (DrawingML).
Requires python-pptx. SVG source directory: _internal/02_svg_source/

Usage:
  python scripts/native_svg_to_ppt.py <project_dir> [-o <output.pptx>]
  (Typically called via: python scripts/pptflow.py <project_dir> export)
"""

import argparse
import base64
import os
import re
import sys
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


CANVAS_PX_W = 1920
CANVAS_PX_H = 1080
PPT_SLIDE_W = Inches(13.333)
PPT_SLIDE_H = Inches(7.5)
PX_TO_EMU = 914400 / CANVAS_PX_W  # 914400 EMU per inch, 13.333in = 1920px


def emu(px_x: float, px_y: float = None) -> int:
    """Convert px to EMU (PowerPoint internal unit)."""
    if px_y is None:
        return int(px_x * PX_TO_EMU)
    return int(px_y * PX_TO_EMU)


def px_to_pt(px: str) -> float:
    """Convert SVG px size to PowerPoint pt (1pt = 1.333px)."""
    val = float(px.replace("px", ""))
    return val * 0.75


def parse_color(color_str: str) -> RGBColor:
    """Parse SVG color string to RGBColor. Handles hex, named, and rgb()."""
    color_str = color_str.strip()
    if color_str.startswith("#"):
        hex_str = color_str[1:]
        if len(hex_str) == 3:
            hex_str = "".join(c * 2 for c in hex_str)
        return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    if color_str.startswith("rgb("):
        m = re.match(r"rgb\((\d+),\s*(\d+),\s*(\d+)\)", color_str)
        if m:
            return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    named = {
        "black": (0, 0, 0), "white": (255, 255, 255), "red": (255, 0, 0),
        "green": (0, 128, 0), "blue": (0, 0, 255), "gray": (128, 128, 128),
        "grey": (128, 128, 128), "transparent": (0, 0, 0),
        "none": (0, 0, 0),
    }
    if color_str.lower() in named:
        r, g, b = named[color_str.lower()]
        return RGBColor(r, g, b)
    return RGBColor(0, 0, 0)


def parse_opacity(elem) -> float:
    """Extract fill-opacity or opacity attribute."""
    op = elem.get("fill-opacity") or elem.get("opacity") or "1"
    try:
        return float(op)
    except ValueError:
        return 1.0


def add_image(slide, href: str, x: float, y: float, w: float, h: float, preserve_aspect: str = None):
    """Add an image to the slide."""
    if href.startswith("data:image"):
        _, encoded = href.split(",", 1)
        img_data = base64.b64decode(encoded)
    else:
        img_path = Path(href)
        if not img_path.exists():
            img_path = Path.cwd() / href
        if not img_path.exists():
            print(f"  [WARN] Image not found: {href}")
            return
        img_data = img_path.read_bytes()

    try:
        slide.shapes.add_picture(
            io.BytesIO(img_data) if "data:" in href else img_data,
            emu(x), emu(y), emu(w), emu(h)
        )
    except Exception as e:
        if "data:" in href:
            import io
            slide.shapes.add_picture(
                io.BytesIO(img_data), emu(x), emu(y), emu(w), emu(h)
            )


def add_rect_shape(slide, x: float, y: float, w: float, h: float,
                   fill_color: str = None, fill_opacity: float = 1.0,
                   stroke_color: str = None, stroke_width: float = 0,
                   rx: float = 0, ry: float = 0):
    """Add a rectangle shape."""
    if fill_color and fill_color.lower() in ("none", "transparent"):
        fill_color = None

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE,
        emu(x), emu(y), emu(w), emu(h)
    )

    if rx > 0:
        adj_val = min(rx, min(w, h) / 2) / min(w, h) * 100000 if min(w, h) > 0 else 0
        try:
            shape.adjustments[0] = adj_val / 100000
        except Exception:
            pass

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = parse_color(fill_color)
        shape.fill.fore_color.brightness = 1.0 - fill_opacity if fill_opacity < 1 else 0
    else:
        shape.fill.background()

    if stroke_color and stroke_color.lower() not in ("none", "transparent"):
        shape.line.color.rgb = parse_color(stroke_color)
        shape.line.width = Pt(stroke_width)
    else:
        shape.line.fill.background()

    return shape


def add_text_box(slide, x: float, y: float, w: float, h: float,
                 text: str, font_family: str = "Calibri",
                 font_size: float = 18, color: str = "#333333",
                 bold: bool = False, align: str = "left"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT, "start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER,
    }.get(align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = text
    run.font.size = Pt(px_to_pt(str(font_size)))
    run.font.color.rgb = parse_color(color)
    run.font.name = font_family
    run.font.bold = bold

    return txBox


def convert_svg_to_slide(slide, svg_root: ET.Element, svg_dir: Path):
    """Convert SVG element tree to PowerPoint slide shapes."""
    ns = svg_root.tag if svg_root.tag.startswith("{") else ""
    ns_uri = ns[1:-1] if ns else ""

    viewbox = svg_root.get("viewBox", "0 0 1920 1080")
    parts = list(map(float, viewbox.split()))
    vb_x, vb_y, vb_w, vb_h = parts if len(parts) == 4 else (0, 0, 1920, 1080)

    def process_element(elem, parent_g_transform=None):
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
        style = elem.get("style", "")

        if tag == "defs":
            return

        if tag == "g":
            for child in elem:
                process_element(child, parent_g_transform)
            return

        x = float(elem.get("x", 0))
        y = float(elem.get("y", 0))
        w = float(elem.get("width", 0))
        h = float(elem.get("height", 0))

        fill = elem.get("fill", style_get(style, "fill", "#333333"))
        stroke = elem.get("stroke", style_get(style, "stroke", "none"))
        stroke_w = float(elem.get("stroke-width", style_get(style, "stroke-width", "0")))
        opacity = parse_opacity(elem)

        if tag == "rect":
            rx = float(elem.get("rx", 0))
            ry = float(elem.get("ry", 0))
            add_rect_shape(slide, x, y, w, h, fill, opacity, stroke, stroke_w, rx, ry)

        elif tag in ("circle", "ellipse"):
            cx = float(elem.get("cx", x))
            cy = float(elem.get("cy", y))
            r = float(elem.get("r", 0)) if tag == "circle" else 0
            rx2 = float(elem.get("rx", r))
            ry2 = float(elem.get("ry", r))

            if fill and fill.lower() not in ("none", "transparent"):
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    emu(cx - rx2), emu(cy - ry2), emu(rx2 * 2), emu(ry2 * 2)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = parse_color(fill)
                if stroke and stroke.lower() not in ("none", "transparent"):
                    shape.line.color.rgb = parse_color(stroke)
                    shape.line.width = Pt(stroke_w)
                else:
                    shape.line.fill.background()

        elif tag == "line":
            x1 = float(elem.get("x1", 0))
            y1 = float(elem.get("y1", 0))
            x2 = float(elem.get("x2", 0))
            y2 = float(elem.get("y2", 0))
            connector = slide.shapes.add_connector(
                1, emu(x1), emu(y1), emu(x2), emu(y2)
            )
            if stroke and stroke.lower() not in ("none", "transparent"):
                connector.line.color.rgb = parse_color(stroke)
                connector.line.width = Pt(stroke_w)

        elif tag in ("path", "polygon", "polyline"):
            if fill and fill.lower() not in ("none", "transparent"):
                pts = elem.get("points", "").strip()
                if not pts and tag == "path":
                    pts = elem.get("d", "")
                if pts:
                    coords = re.findall(r"[-]?\d+\.?\d*", pts)
                    if len(coords) >= 4:
                        x_vals = [float(coords[i]) for i in range(0, len(coords), 2)]
                        y_vals = [float(coords[i]) for i in range(1, len(coords), 2)]
                        min_x, max_x = min(x_vals), max(x_vals)
                        min_y, max_y = min(y_vals), max(y_vals)
                        pw = max_x - min_x if max_x > min_x else 10
                        ph = max_y - min_y if max_y > min_y else 10
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE if tag in ("polygon", "polyline") else MSO_SHAPE.FREEHOLD,
                            emu(min_x), emu(min_y), emu(pw), emu(ph)
                        )
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = parse_color(fill)
                        if stroke and stroke.lower() not in ("none", "transparent"):
                            shape.line.color.rgb = parse_color(stroke)
                            shape.line.width = Pt(stroke_w)
                        else:
                            shape.line.fill.background()

        elif tag == "text":
            font_family = elem.get("font-family", "Microsoft YaHei")
            font_size = float(elem.get("font-size", "24").replace("px", ""))
            color = elem.get("fill", "#333333")
            bold = elem.get("font-weight", "normal") in ("bold", "600", "700", "800")
            text_anchor = elem.get("text-anchor", "start")

            text_content = "".join(elem.itertext())
            if not text_content.strip():
                return

            add_text_box(slide, x, y, w if w else 400, h if h else 60,
                         text_content, font_family, font_size, color, bold,
                         text_anchor)

        elif tag == "image":
            href = elem.get("href") or elem.get("{http://www.w3.org/1999/xlink}href", "")
            if href:
                pa = elem.get("preserveAspectRatio", "xMidYMid slice")
                add_image(slide, href, x, y, w, h, pa)

        elif tag == "use":
            href = elem.get("href") or elem.get("{http://www.w3.org/1999/xlink}href", "")
            if href and href.startswith("data:icon"):
                icon_size = float(elem.get("width", 48))
                icon_fill = elem.get("fill", fill)
                icon_data = href.split(",", 1)[1] if "," in href else ""
                if icon_data:
                    import io
                    try:
                        icon_bytes = base64.b64decode(icon_data)
                        slide.shapes.add_picture(io.BytesIO(icon_bytes), emu(x), emu(y),
                                                  emu(icon_size), emu(icon_size))
                    except Exception:
                        pass

    for child in svg_root:
        process_element(child)


def style_get(style: str, key: str, default: str = None) -> str:
    """Extract a value from inline CSS style string."""
    if not style:
        return default
    for part in style.split(";"):
        part = part.strip()
        if part.startswith(key + ":"):
            return part.split(":", 1)[1].strip()
    return default


def convert_project(project_dir: str, output_path: str = None):
    """Convert all SVG pages in a project to a single PPTX."""
    p = Path(project_dir)
    svg_dir = p / "_internal" / "02_svg_source"
    output_path = output_path or str(p / "final_deck.pptx")

    if not svg_dir.exists():
        print(f"[ERROR] SVG source directory not found: {svg_dir}", file=sys.stderr)
        # Try project root svg_output/
        svg_dir = p / "svg_output"
        if not svg_dir.exists():
            print(f"[ERROR] No svg_output/ directory either", file=sys.stderr)
            sys.exit(1)

    approval = os.environ.get("SMART_SVG_EXPORT_APPROVED_BY_PPTFLOW", "")
    if approval != "1":
        print("[ERROR] Export not approved by pptflow.py.", file=sys.stderr)
        print("[ERROR] Run: python scripts/pptflow.py <project_dir> export", file=sys.stderr)
        sys.exit(1)

    print(f"[CONVERT] Reading SVGs from {svg_dir}")

    svg_files = sorted(svg_dir.glob("*.svg"))
    if not svg_files:
        print(f"[ERROR] No SVG files found", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = PPT_SLIDE_W
    prs.slide_height = PPT_SLIDE_H
    blank_layout = prs.slide_layouts[6]

    for svg_path in svg_files:
        print(f"  [SLIDE] {svg_path.name}")
        slide = prs.slides.add_slide(blank_layout)

        try:
            tree = ET.parse(str(svg_path))
            root = tree.getroot()
            convert_svg_to_slide(slide, root, svg_dir)
        except ET.ParseError as e:
            print(f"  [ERROR] XML parse error in {svg_path.name}: {e}", file=sys.stderr)
            continue

    prs.save(str(output_path))
    print(f"\n[CONVERT] [OK] PPTX saved to {output_path}")


def main():
    parser = argparse.ArgumentParser(description="SVG to PPTX Converter")
    parser.add_argument("project_dir", help="Project directory")
    parser.add_argument("-o", "--output", help="Output PPTX path")
    args = parser.parse_args()
    convert_project(args.project_dir, args.output)


if __name__ == "__main__":
    import io
    main()
