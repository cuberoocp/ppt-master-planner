#!/usr/bin/env python3
"""Unified source to Markdown converter.

Converts various source formats to Markdown for downstream processing.
Supports PDF, DOCX, XLSX, PPTX, and URLs.

Usage:
  python scripts/source_to_md/pdf_to_md.py <input>
  python scripts/source_to_md/doc_to_md.py <input>
  python scripts/source_to_md/excel_to_md.py <input>
  python scripts/source_to_md/ppt_to_md.py <input>
  python scripts/source_to_md/web_to_md.py <URL>
"""

import argparse
import os
import re
import sys
from pathlib import Path


def convert_pdf(input_path: str, output_path: str = None):
    """Extract text from PDF to Markdown."""
    try:
        from pypdf import PdfReader
    except ImportError:
        print("[ERROR] pypdf not installed. Run: pip install pypdf", file=sys.stderr)
        sys.exit(1)

    reader = PdfReader(input_path)
    md_lines = [f"# PDF: {Path(input_path).name}\n"]
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        md_lines.append(f"\n## Page {i + 1}\n\n{text}")

    output = output_path or str(Path(input_path).with_suffix(".md"))
    with open(output, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))
    print(f"[OK] PDF -> Markdown: {output}")
    return output


def convert_docx(input_path: str, output_path: str = None):
    """Extract text from DOCX to Markdown."""
    try:
        from docx import Document
    except ImportError:
        print("[ERROR] python-docx not installed. Run: pip install python-docx", file=sys.stderr)
        sys.exit(1)

    doc = Document(input_path)
    md_lines = [f"# Document: {Path(input_path).name}\n"]
    for para in doc.paragraphs:
        style = para.style.name.lower() if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if "heading 1" in style:
            md_lines.append(f"\n# {text}\n")
        elif "heading 2" in style:
            md_lines.append(f"\n## {text}\n")
        elif "heading 3" in style:
            md_lines.append(f"\n### {text}\n")
        else:
            md_lines.append(text)

    output = output_path or str(Path(input_path).with_suffix(".md"))
    with open(output, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))
    print(f"[OK] DOCX -> Markdown: {output}")
    return output


def convert_excel(input_path: str, output_path: str = None):
    """Extract data from XLSX/XLSM to Markdown tables."""
    try:
        import openpyxl
    except ImportError:
        print("[ERROR] openpyxl not installed. Run: pip install openpyxl", file=sys.stderr)
        sys.exit(1)

    wb = openpyxl.load_workbook(input_path, read_only=True, data_only=True)
    md_lines = [f"# Excel: {Path(input_path).name}\n"]
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        md_lines.append(f"\n## Sheet: {sheet_name}\n")
        rows = list(ws.iter_rows(values_only=True))
        if rows:
            header = "| " + " | ".join(str(c or "") for c in rows[0]) + " |"
            sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
            md_lines.append(header)
            md_lines.append(sep)
            for row in rows[1:]:
                md_lines.append("| " + " | ".join(str(c or "") for c in row) + " |")

    output = output_path or str(Path(input_path).with_suffix(".md"))
    with open(output, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))
    print(f"[OK] Excel -> Markdown: {output}")
    return output


def convert_pptx(input_path: str, output_path: str = None):
    """Extract text from PPTX to Markdown."""
    try:
        from pptx import Presentation
    except ImportError:
        print("[ERROR] python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(input_path)
    md_lines = [f"# PowerPoint: {Path(input_path).name}\n"]
    for i, slide in enumerate(prs.slides):
        md_lines.append(f"\n## Slide {i + 1}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        if para.level > 0:
                            md_lines.append("  " * para.level + "- " + text)
                        else:
                            md_lines.append(text)

    output = output_path or str(Path(input_path).with_suffix(".md"))
    with open(output, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))
    print(f"[OK] PPTX -> Markdown: {output}")
    return output


def convert_web(url: str, output_path: str = None):
    """Fetch URL, download images, and convert to Markdown with local image refs."""
    try:
        import requests
        from bs4 import BeautifulSoup
    except ImportError:
        print("[ERROR] requests/beautifulsoup4 not installed.", file=sys.stderr)
        sys.exit(1)

    try:
        resp = requests.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
    except Exception as e:
        print(f"[ERROR] Failed to fetch URL: {e}", file=sys.stderr)
        sys.exit(1)

    soup = BeautifulSoup(resp.text, "html.parser")
    title = soup.title.string if soup.title else Path(url).name
    md_lines = [f"# Web: {title}\n", f"> Source: {url}\n"]

    # Determine output paths
    safe_name = re.sub(r'[\\/*?:\"<>|]', '_', url.split('?')[0].rstrip('/').split('/')[-1] or 'webpage')
    output = output_path or f"web_{safe_name}.md"
    output_dir = Path(output).parent
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    # Track downloaded images to avoid duplicates
    downloaded = {}

    # Iterate body children in order to preserve text + image placement
    body = soup.find("body") or soup
    for tag in body.descendants:
        if tag.name == "img":
            src = tag.get("src", "")
            if not src or src.startswith("data:"):
                continue
            if not src.startswith(("http://", "https://")):
                from urllib.parse import urljoin
                src = urljoin(url, src)

            if src in downloaded:
                local_path = downloaded[src]
            else:
                try:
                    img_resp = requests.get(src, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
                    img_resp.raise_for_status()
                    ext = Path(src.split("?")[0]).suffix or ".png"
                    fname = f"img_{len(downloaded) + 1}{ext}"
                    local_path = images_dir / fname
                    local_path.write_bytes(img_resp.content)
                    downloaded[src] = local_path
                    print(f"  [IMG] Saved {fname} ({len(img_resp.content)} bytes)")
                except Exception as e:
                    print(f"  [IMG] WARN: failed to download {src}: {e}")
                    continue

            alt = tag.get("alt", "")
            md_lines.append(f"\n![{alt}](images/{local_path.name})\n")

        elif tag.name in ("h1", "h2", "h3", "h4"):
            text = tag.get_text(strip=True)
            if text:
                level = "#" * int(tag.name[1])
                md_lines.append(f"\n{level} {text}")

        elif tag.name == "p":
            text = tag.get_text(strip=True)
            if text:
                md_lines.append(f"\n{text}")

        elif tag.name == "li":
            text = tag.get_text(strip=True)
            if text:
                md_lines.append(f"- {text}")

        elif tag.name == "table":
            rows = tag.find_all("tr")
            for row in rows:
                cells = row.find_all(["td", "th"])
                md_lines.append("| " + " | ".join(c.get_text(strip=True) for c in cells) + " |")

    with open(output, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))
    print(f"[OK] Web -> Markdown: {output} ({len(downloaded)} images saved)")
    return output


def main():
    script_name = Path(__file__).name
    print(f"Usage: Use the specific script for your source type:")
    print(f"  python scripts/source_to_md/pdf_to_md.py <PDF_file>")
    print(f"  python scripts/source_to_md/doc_to_md.py <DOCX_file>")
    print(f"  python scripts/source_to_md/excel_to_md.py <XLSX_file>")
    print(f"  python scripts/source_to_md/ppt_to_md.py <PPTX_file>")
    print(f"  python scripts/source_to_md/web_to_md.py <URL>")


if __name__ == "__main__":
    main()
